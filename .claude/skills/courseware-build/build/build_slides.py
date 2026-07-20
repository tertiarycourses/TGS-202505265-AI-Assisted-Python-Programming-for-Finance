#!/usr/bin/env python3
"""Generate the AI Assisted Python Programming for Finance course slide deck (all-white Tertiary house style).

Design helpers are the same set used by the tertiary-course-slides skill that
produced the n8n reference deck (cover, section, content, two_col, cards3,
big_statement, step_slide, test_slide, brk). Content is driven entirely by
course_data.py + data_domainN.py so the deck stays 100% aligned with the LP,
LG and labs.
"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C
from data_domain1 import DOMAIN1
from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3
from data_domain4 import DOMAIN4
from data_domain5 import DOMAIN5
from data_domain6 import DOMAIN6
from data_domain7 import DOMAIN7
from data_domain8 import DOMAIN8
from data_domain9 import DOMAIN9
ACTIVITIES = (DOMAIN1 + DOMAIN2 + DOMAIN3 + DOMAIN4 + DOMAIN5
              + DOMAIN6 + DOMAIN7 + DOMAIN8 + DOMAIN9)

def _find_repo(start):
    """Locate the course repo (a dir containing both courseware/ and labs/).
    Env COURSE_REPO overrides. Keeps the build working wherever the skill lives."""
    env = os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env):
        return env
    d = start
    for _ in range(8):
        d = os.path.dirname(d)
        if os.path.isdir(os.path.join(d, "courseware")) and os.path.isdir(os.path.join(d, "labs")):
            return d
    return os.path.dirname(os.path.dirname(HERE))
REPO = _find_repo(HERE)
ASSETS = os.path.join(os.path.dirname(HERE), "assets")   # co-located with the skill

# ---------------- palette (matches reference) ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Arial"
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=color if lvl==0 else GREY
        r.font.name="Arial"; r.font.bold=(lvl==0 and isinstance(it,tuple) and len(it)>2 and it[2])
    return tb

PAGE={"n":0}
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}",9,GREY,False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,GREY,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,GREY,False)]],align=PP_ALIGN.RIGHT)
def head(s,title,kicker=None,kcolor=BLUE):
    """Header block. The title font shrinks and the divider rule moves down so a
    two-line title is never struck through by the rule."""
    rect(s,0,0,SW,SH,WHITE)
    TW=11.75
    pt=_fit_pt(title,TW,29,20,2)
    nlines=_est_lines(title,TW,pt,bold=True)
    th=Inches(0.52*nlines+0.16)
    rule_y=Inches(0.9)+th+Inches(0.06)
    accent_h=max(Inches(1.55),rule_y-Inches(0.05))
    rect(s,0,0,Inches(0.28),accent_h,kcolor)
    if kicker: txt(s,Inches(0.85),Inches(0.5),Inches(11.6),Inches(0.4),[[(kicker,14,kcolor,True)]])
    txt(s,Inches(0.85),Inches(0.9),Inches(TW),th,[[(title,pt,INK,True)]])
    rect(s,Inches(0.85),rule_y,Inches(11.63),Inches(0.02),LINE)
    s._body_top=rule_y+Inches(0.25)
    return s
def _logo(name):
    p=os.path.join(ASSETS,name)
    return p if os.path.exists(p) else None


# ---------------- text measurement / autofit ----------------
# Titles and body panels must be sized from the ACTUAL text, not a fixed line
# count — otherwise anything that wraps silently collides with the element below.
def _est_lines(text, width_in, pt, bold=False):
    """Estimate wrapped line count for Arial at `pt` inside `width_in` inches."""
    if not text: return 1
    # average Arial glyph advance ~0.50 em regular, ~0.53 em bold
    char_w_in = (pt / 72.0) * (0.53 if bold else 0.50)
    per_line = max(8, int(width_in / char_w_in))
    lines = 1; cur = 0
    for word in str(text).split():
        add = len(word) + (1 if cur else 0)
        if cur + add > per_line:
            lines += 1; cur = len(word)
        else:
            cur += add
    return lines

def _fit_pt(text, width_in, max_pt, min_pt, max_lines):
    """Largest font size (<=max_pt) whose wrapped text fits in max_lines."""
    pt = max_pt
    while pt > min_pt and _est_lines(text, width_in, pt, bold=True) > max_lines:
        pt -= 1
    return pt

def _autofit(tb, shrink_to=0.6):
    """Let PowerPoint shrink text that still overflows its frame."""
    from pptx.oxml.ns import qn
    bodyPr = tb.text_frame._txBody.find(qn('a:bodyPr'))
    if bodyPr is None: return
    for tag in ('a:noAutofit','a:spAutoFit','a:normAutofit'):
        el = bodyPr.find(qn(tag))
        if el is not None: bodyPr.remove(el)
    el = bodyPr.makeelement(qn('a:normAutofit'), {})
    el.set('fontScale', str(int(shrink_to*100000)))
    el.set('lnSpcReduction', '10000')
    bodyPr.append(el)
    return tb

# ---------------- slide templates ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    # course badge (top-right) — IBF accreditation badge, else text fallback
    badge=_logo("ibf-badge.png")
    if badge:
        s.shapes.add_picture(badge,Inches(10.05),Inches(0.6),width=Inches(2.5))
    else:
        rect(s,Inches(10.85),Inches(0.72),Inches(1.7),Inches(1.0),BLUE)
        txt(s,Inches(10.85),Inches(0.82),Inches(1.7),Inches(0.5),[[("IBF-STS",17,WHITE,True)]],align=PP_ALIGN.CENTER)
        txt(s,Inches(10.85),Inches(1.26),Inches(1.7),Inches(0.4),[[("ACCREDITED",8,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES  ·  IBF-STS ACCREDITED",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[(C.TITLE,40,INK,True)]])
    rect(s,Inches(0.92),Inches(4.75),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(5.05),Inches(12),Inches(1.4),
        [[(f"Course Code: {C.COURSE_CODE}   ·   IBF Standard: {C.ACCREDITATION['standard']}",16,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.5),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,GREY,False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])

def section(kicker,title,n,sub=""):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,BLUE)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,BLUE,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.55),Inches(11),Inches(0.8),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,RGBColor(0xE2,0xE8,0xF0),True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def _top(s,default=1.95):
    return getattr(s,"_body_top",Inches(default))
def content(title,items,kicker=None,size=20):
    s=head(slide(),title,kicker); y=_top(s)
    tb=bullets(s,Inches(0.85),y,Inches(11.6),int(Inches(6.78)-y),items,size=size)
    _autofit(tb,0.62); footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead="",bottom=None):
    s=head(slide(),title,kicker); y=_top(s)
    bot=bottom if bottom is not None else Inches(6.78)
    ph=int(bot-y)
    rect(s,Inches(0.85),y,Inches(5.7),ph,LIGHT); rect(s,Inches(6.95),y,Inches(5.55),ph,LIGHT)
    hy=y+Inches(0.2); by=y+Inches(0.67); bh=int(ph-Inches(0.82))
    if lhead: txt(s,Inches(1.1),hy,Inches(5.2),Inches(0.4),[[(lhead,16,BLUE,True)]])
    if rhead: txt(s,Inches(7.2),hy,Inches(5.0),Inches(0.4),[[(rhead,16,TEAL,True)]])
    _autofit(bullets(s,Inches(1.1),by,Inches(5.2),bh,left,size=13,gap=5),0.6)
    _autofit(bullets(s,Inches(7.2),by,Inches(5.05),bh,right,size=13,gap=5,mcolor=TEAL),0.6)
    footer(s); return s
def cards3(title,cards,kicker):
    s=head(slide(),title,kicker); y=_top(s); xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    ph=int(Inches(6.78)-y)
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,y,Inches(3.65),ph,LIGHT); rect(s,x,y,Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),y+Inches(0.25),Inches(3.2),Inches(0.6),[[(c[1],19,col,True)]])
        _autofit(bullets(s,x+Inches(0.25),y+Inches(1.0),Inches(3.2),int(ph-Inches(1.15)),
                         c[2],size=14,mcolor=col,gap=9),0.6)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
import math
PALETTE=[BLUE,TEAL,VIOLET,AMBER]
def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    """Grid of light panels, each with a coloured icon/number badge + text.
    items: list of strings (or (title,caption) tuples). Much richer than a bullet list."""
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=_top(s); TOTW=Inches(11.63); AREAH=int(Inches(6.78)-Y0)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,19,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],size+2,INK,True)],[(it[1],size-2,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,size,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def flow_h(title,steps,kicker=None,color=BLUE):
    """Horizontal numbered flow: coloured chips connected by chevrons."""
    s=head(slide(),title,kicker,kcolor=color)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.55); ch=Inches(3.15); bd=Inches(0.82)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),color)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,[[(str(i+1),30,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.16),int(y+Inches(1.55)),cw-Inches(0.32),int(ch-Inches(1.7)),[[(st,14,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,color,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    """Profile-card layout: avatar badge + name/role panel on the left, labelled
    info tiles on the right. rows: list of (LABEL, value); blank value → fill-in line."""
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    y0=_top(s); ph=int(Inches(6.78)-y0)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,y0,lw,ph,LIGHT); rect(s,lx,y0,lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,GREY,False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=y0; rh=ph
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,14,INK,False)] if val else [("____________________________________________",13,LINE,False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def activity_overview(tag,title,desc,build,services,kicker):
    s=head(slide(),title,kicker,kcolor=TEAL); y=_top(s)-Inches(0.1)
    rect(s,Inches(0.85),y,Inches(1.7),Inches(0.5),TEAL)
    txt(s,Inches(0.85),y+Inches(0.05),Inches(1.7),Inches(0.4),[[(tag,16,WHITE,True)]],align=PP_ALIGN.CENTER)
    dy=y+Inches(0.7)
    dpt=18 if _est_lines(desc,11.7,18)<=3 else 15
    dlines=_est_lines(desc,11.7,dpt)
    dh=Inches(0.30*dlines+0.12)
    _autofit(txt(s,Inches(0.85),dy,Inches(11.7),dh,[[(desc,dpt,INK,False)]]),0.7)
    # the build panel takes whatever height is left, and never rides up over the text
    by=dy+dh+Inches(0.2); bh=int(Inches(6.78)-by)
    if bh<Inches(1.3): bh=Inches(1.3)
    rect(s,Inches(0.85),by,Inches(11.7),bh,LIGHT)
    txt(s,Inches(1.1),by+Inches(0.18),Inches(11),Inches(0.34),[[("You'll build",13,BLUE,True)]])
    _autofit(txt(s,Inches(1.1),by+Inches(0.56),Inches(11.2),int(bh-Inches(1.15)),
                 [[(build,16,INK,True)]]),0.7)
    txt(s,Inches(1.1),int(by+bh-Inches(0.5)),Inches(11.2),Inches(0.4),
        [[("Tools:  ",12,GREY,True),(services,12,GREY,False)]])
    footer(s); return s
def step_slide(kicker,act_title,n,total,text,cmd=""):
    s=head(slide(),act_title,kicker,TEAL)
    oval(s,Inches(0.85),Inches(2.5),Inches(1.4),Inches(1.4),TEAL)
    txt(s,Inches(0.85),Inches(2.74),Inches(1.4),Inches(0.9),[[(str(n),38,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.95),Inches(1.95),Inches(11),Inches(0.4),[[(f"STEP {n} OF {total}",13,GREY,True)]])
    txt(s,Inches(2.55),Inches(2.4),Inches(10.1),Inches(1.3),[[(text,23,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if cmd:
        rect(s,Inches(2.55),Inches(4.15),Inches(10.1),Inches(0.95),RGBColor(0x0B,0x12,0x20))
        txt(s,Inches(2.8),Inches(4.28),Inches(9.7),Inches(0.7),[[("$ "+cmd,13,RGBColor(0x9C,0xDC,0xFE),False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def test_slide(act_title,text,kicker):
    s=head(slide(),act_title,kicker,TEAL)
    rect(s,Inches(0.85),Inches(2.3),Inches(11.7),Inches(2.6),RGBColor(0xE8,0xF7,0xEE))
    txt(s,Inches(1.2),Inches(2.6),Inches(11),Inches(0.5),[[("✅  Test it",20,RGBColor(0x12,0x7A,0x3E),True)]])
    txt(s,Inches(1.2),Inches(3.3),Inches(11),Inches(1.4),[[(text,18,INK,False)]]); footer(s); return s
def brk(kind,dur,color=AMBER):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1

def steps_slide(kicker,act_title,steps,part=None,parts=1):
    """Compact numbered step list. Row heights are PROPORTIONAL TO MEASURED TEXT,
    so a long AI-prompt step gets the space it needs and never overprints its
    neighbours."""
    title=act_title if parts==1 else f"{act_title}  ({part}/{parts})"
    s=head(slide(),title,kicker,TEAL)
    Y0=_top(s); AREA=int(Inches(6.78)-Y0)
    gy=Inches(0.10); n=len(steps)
    avail=AREA-gy*(n-1)
    # weight each row by how many lines its text needs at the base size
    TXT_W_CMD, TXT_W_PLAIN = 6.35, 10.6
    weights=[]
    for instr,cmd in steps:
        w=TXT_W_CMD if cmd else TXT_W_PLAIN
        weights.append(max(1,_est_lines(instr,w,12)))
    tot=sum(weights)
    heights=[max(int(Inches(0.42)),int(avail*wt/tot)) for wt in weights]
    over=sum(heights)-avail
    while over>0:                      # trim proportionally if we overshot
        for i in range(len(heights)):
            if heights[i]>int(Inches(0.42)) and over>0:
                heights[i]-=1; over-=1
    y=Y0
    for i,((instr,cmd),rh) in enumerate(zip(steps,heights)):
        rect(s,Inches(0.85),y,Inches(11.63),rh,LIGHT)
        bd=min(Inches(0.46),max(Inches(0.26),rh-Inches(0.1)))
        oval(s,Inches(1.0),int(y+rh/2-bd/2),bd,bd,TEAL)
        txt(s,Inches(1.0),int(y+rh/2-bd/2),bd,bd,[[(str(i+1),13,WHITE,True)]],
            align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=Inches(1.62)
        if cmd:
            _autofit(txt(s,tx,y+Inches(0.04),Inches(TXT_W_CMD),int(rh-Inches(0.08)),
                         [[(instr,12,INK,False)]],anchor=MSO_ANCHOR.MIDDLE),0.62)
            cw=Inches(4.25); ch=min(int(rh-Inches(0.12)),int(Inches(0.62)))
            cy=int(y+rh/2-ch/2)
            rect(s,Inches(8.15),cy,cw,ch,RGBColor(0x0B,0x12,0x20))
            _autofit(txt(s,Inches(8.28),cy,Inches(3.99),ch,
                         [[("$ "+cmd,10,RGBColor(0x9C,0xDC,0xFE),False)]],
                         anchor=MSO_ANCHOR.MIDDLE),0.55)
        else:
            _autofit(txt(s,tx,y+Inches(0.04),Inches(TXT_W_PLAIN),int(rh-Inches(0.08)),
                         [[(instr,12,INK,False)]],anchor=MSO_ANCHOR.MIDDLE),0.62)
        y+=rh+gy
    footer(s); return s

def lab_steps(kicker,act_title,steps,per=6):
    """Chunk a lab's steps across slides, using a LINE BUDGET rather than a fixed
    count so slides holding long prompts carry fewer steps."""
    BUDGET=17          # ~17 text lines fit the content area comfortably
    chunks=[]; cur=[]; used=0
    for instr,cmd in steps:
        w=6.35 if cmd else 10.6
        need=max(1,_est_lines(instr,w,12))
        if cur and (used+need>BUDGET or len(cur)>=per):
            chunks.append(cur); cur=[]; used=0
        cur.append((instr,cmd)); used+=need
    if cur: chunks.append(cur)
    for ci,ch in enumerate(chunks,1):
        steps_slide(kicker,act_title,ch,part=ci,parts=len(chunks))

# ============================================================ BUILD
cover()

def admin_attendance():
    content("Digital Attendance (Mandatory)",[
     "It is mandatory to take the AM, PM and Assessment digital attendance for this funded course.",
     "The trainer/administrator displays the digital attendance QR code from the training portal.",
     "Scan the QR code with your mobile phone camera and submit your attendance.",
     "A minimum of 75% attendance is required to be eligible for assessment and IBF-STS funding."],
     kicker="TRAQOM \u00b7 DIGITAL ATTENDANCE")

def assessment_flow():
    flow_h("Assessment Flow",[
     "TRAQOM survey \u2014 scan the QR code on the LMS",
     "Assessment digital attendance \u2014 scan the QR code",
     "Sit WA (SAQ) then the Practical Test \u2014 open book",
     "Submit your answers and code on the LMS",
     "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")

# ---------------- ADMIN ----------------
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
admin_attendance()
trainer_slide("YOUR TRAINER \u00b7 GENERAL","Your Trainer","General Trainer template \u2014\nto be completed by the trainer",
 [("Name",""),("Title / Designation",""),("Qualifications",""),
  ("Areas of expertise",""),("Training & industry experience",""),("Contact","")],
 initials="?",accent=GREY)
trainer_slide("YOUR TRAINER",C.TRAINER,"Principal Trainer\nTertiary Infotech Academy Pte. Ltd.",
 [("Role","Principal Trainer, Tertiary Infotech Academy Pte. Ltd."),
  ("Expertise","Python programming, data analytics and AI-assisted software development for financial services."),
  ("Delivers","IBF-accredited and WSQ courses on Python, data analytics, machine learning and AI."),
  ("Founder","Founder and lead instructor at Tertiary Infotech / Tertiary Courses.")],
 initials="AA",accent=BLUE)
content("Let's Know Each Other",[
 "Your name, organisation and role in financial services.",
 "Your experience with Python, Excel or data analysis (if any).",
 "The financial analysis or reporting task you most want to automate after this course."],kicker="ICE-BREAKER")
tile_grid("Ground Rules",[
 "Set your mobile phone to silent mode.","Participate actively \u2014 no question is too small.",
 "Mutual respect: agree to disagree.","One conversation at a time.",
 "Be punctual; return from breaks on time.","75% attendance is required."],
 kicker="HOUSEKEEPING",cols=2,size=15)
content("Download Course Material",[
 "Access your slides, Learner Guide, lab files and assessment on the LMS/TMS portal.",
 "Portal: https://lms-tms.tertiaryinfotech.com",
 "Download the slides and Learner Guide \u2014 you may use them during the open-book assessment.",
 "The lab source files and the combined Jupyter notebook are in the Activities folder."],kicker="COURSE PORTAL \u00b7 LMS/TMS")
tile_grid("About IBF-STS",[
 ("Accreditation",f"Accredited under the {C.ACCREDITATION['framework']}, IBF Standard {C.ACCREDITATION['standard']}."),
 ("Funding","Eligible for the IBF Standards Training Scheme (IBF-STS), subject to all eligibility criteria being met."),
 ("Subsidy","IBF-STS provides 50%\u201370% course fee subsidy for direct training costs, capped at S$3,000 per candidate per course."),
 ("Suitability","Participants should assess the suitability of the course and its relevance to their business activities or job roles."),
 ("Find out more",f"Visit {C.ACCREDITATION['website']} for the full eligibility criteria and quantum.")],
 kicker="ACCREDITATION & FUNDING",cols=1,size=14)

# Lesson plan overview \u2014 4 days, derived from DAY_TOPICS
def _day_items(d):
    out=[(f"Day {d} \u2014 {C.DAY_THEMES[d]}",0)]
    for tn in C.DAY_TOPICS[d]:
        t=next(x for x in C.TOPICS if x["num"]==tn)
        la=[a for a in ACTIVITIES if a["topic"]==tn]
        rng=f" (Labs {la[0]['num']}\u2013{la[-1]['num']})" if la else ""
        out.append((f"Topic {tn}: {t['title']}{rng}",1))
    return out
_lp=two_col("Lesson Plan \u2014 4 Days, 8 hours/day",
 _day_items(1)+_day_items(2),
 _day_items(3)+_day_items(4),
 kicker="SCHEDULE",lhead="Days 1\u20132",rhead="Days 3\u20134")
txt(_lp,Inches(0.85),Inches(6.72),Inches(11.63),Inches(0.3),
    [[("Daily timing:  ",12,BLUE,True),
      ("9:30am\u20136:30pm  \u00b7  1-hour lunch break  \u00b7  tea breaks counted within training time",12,GREY,False)]])

def _lo_short(lo,limit=118):
    body=lo.split(":",1)[1].strip()
    return body if len(body)<=limit else body[:limit].rsplit(" ",1)[0]+"\u2026"
tile_grid("Learning Outcomes",[(lo.split(":",1)[0], _lo_short(lo)) for lo in C.LEARNING_OUTCOMES],
 kicker="WHAT YOU'LL ACHIEVE",cols=2,size=11)
def _short_sub(sub,maxparts=4):
    parts=[x.strip() for x in sub.split("\u00b7")]
    out=" \u00b7 ".join(parts[:maxparts])
    return out+(" \u00b7 \u2026" if len(parts)>maxparts else "")
tile_grid("Agenda \u2014 The Nine Topics",
 [(f"Topic {t['num']}: {t['title']}", _short_sub(t["subtitle"])) for t in C.TOPICS],
 kicker="COURSE OUTLINE",cols=2,size=11)
content("Briefing for Assessment",[
 "Place phones and other materials under the table or on the floor.",
 "No photos or recording of assessment scripts.","No discussion during the assessment.",
 "Use a black/blue pen for hard-copy assessments.","No liquid paper / correction tape.",
 "Scripts are collected when time is up."],kicker="BEFORE YOU START")
content("Assessment",[
 C.ASSESSMENT["written"], C.ASSESSMENT["practical"],
 f"Passing score: {C.ASSESSMENT['passing_score']}.",
 "Format: Open Book \u2014 slides, Learner Guide and lab files only.",
 C.ASSESSMENT["note"],"An appeal process is available if required."],kicker="FINAL ASSESSMENT")
assessment_flow()

# ---------------- CORE CONCEPTS ----------------
section("CORE CONCEPTS","Python and AI in Financial Services","")
tile_grid("Why Python for Finance?",[
 ("The #1 language","Python is consistently ranked the leading programming language, and is the default choice across financial services."),
 ("Readable and fast to write","Clear syntax means a financial analyst can read, audit and maintain the code \u2014 not just the developer who wrote it."),
 ("A vast finance ecosystem","pandas, NumPy, matplotlib, yfinance, scikit-learn and Streamlit cover the entire analysis-to-application pipeline."),
 ("Free and open source","No licence cost, and it runs everywhere \u2014 laptop, server, or a browser tab via Google Colab.")],
 kicker="OVERVIEW",cols=2,size=14)
tile_grid("Use Cases of Python in Financial Services",[
 ("Algorithmic Trading","Developing and backtesting algorithms for stocks, forex and commodities."),
 ("Financial Analysis & Modelling","pandas and NumPy for analysis; matplotlib and seaborn for visualisation."),
 ("Portfolio Management","Construction, analysis and optimisation of portfolios against risk and return targets."),
 ("Fraud Detection","Machine-learning libraries such as scikit-learn and TensorFlow to detect anomalous transactions."),
 ("Risk Management","Quantifying credit, market and operational risk, and stress-testing the results."),
 ("Reporting & Automation","Replacing manual spreadsheet work with reproducible, auditable scripts.")],
 kicker="WHERE IT IS USED",cols=2,size=13)
cards3("From Business Requirement to Working Code",[
 (BLUE,"1. Understand",["Clarify the business requirement","Identify inputs, outputs and rules","Define what 'correct' means"]),
 (TEAL,"2. Build",["Break it into functions","Prompt an AI tool for a first draft","Review, test and correct the code"]),
 (VIOLET,"3. Deliver",["Handle errors and edge cases","Document for business users","Publish as a report or Streamlit app"])],kicker="METHODOLOGY")
big_statement("AI writes the first draft. You are accountable for the result.",
 "Generated code must be read, tested and understood before it touches a financial system \u2014 a plausible-looking error in a risk model is worse than no model at all.",
 "THE GOLDEN RULE",color=AMBER)
tile_grid("Your AI-Assisted Toolkit",[
 ("Google Colab","Browser-based notebooks with a built-in Gemini assistant \u2014 nothing to install, ideal for the classroom."),
 ("Visual Studio Code","The industry-standard editor, with GitHub Copilot for inline AI completion."),
 ("Cursor","An AI-first editor that edits across a whole project from a natural-language instruction."),
 ("AntiGravity","Google's agentic development environment for larger, multi-step builds.")],
 kicker="TOOLS YOU WILL USE",cols=2,size=14)
tile_grid("Managing Python Projects with uv",[
 ("What uv is","A single, very fast tool that replaces pip, virtualenv and pyenv for creating and managing Python projects."),
 ("uv init","Creates a new project folder with a pyproject.toml and an isolated environment."),
 ("uv add","Installs a dependency and records it, so the project is reproducible on any machine."),
 ("uv run","Runs your script inside the project environment \u2014 no manual activation step."),
 ("Why it matters","Reproducibility is an audit requirement in finance: the same code and the same versions must give the same numbers.")],
 kicker="EVERY LAB USES uv",cols=1,size=14,accent=TEAL)
flow_h("The Lab Workflow",[
 "uv init \u2014 create the project",
 "uv add \u2014 install pandas, yfinance, streamlit",
 "Prompt the AI tool for a first draft",
 "Review, correct and run with uv run",
 "Test against the expected result"],kicker="HOW EVERY LAB RUNS",color=TEAL)

# ---------------- TOPICS + ACTIVITIES ----------------
TOPIC_ACTS = {t["num"]: [a for a in ACTIVITIES if a["topic"]==t["num"]] for t in C.TOPICS}
CARD_COLORS=[BLUE,TEAL,VIOLET]
DAY_OF_TOPIC={tn:d for d,tns in C.DAY_TOPICS.items() for tn in tns}
LAST_TOPIC_OF_DAY={tns[-1]:d for d,tns in C.DAY_TOPICS.items()}
FIRST_TOPIC_OF_DAY={tns[0]:d for d,tns in C.DAY_TOPICS.items()}

for t in C.TOPICS:
    tn=t["num"]
    if tn in FIRST_TOPIC_OF_DAY and FIRST_TOPIC_OF_DAY[tn]>1:
        d=FIRST_TOPIC_OF_DAY[tn]
        big_statement(f"Day {d}",C.DAY_THEMES[d],f"START OF DAY {d}",color=BLUE)
        admin_attendance()
    section(f"TOPIC {t['code']}", t["title"], t["code"], t["subtitle"])
    tile_grid(f"Key Concepts \u2014 {t['title']}", t["concepts"],
              kicker=f"TOPIC WEIGHTING {t['weighting']}", cols=2, size=14)
    acts=TOPIC_ACTS[tn]
    third=(len(acts)+2)//3
    groups=[acts[i:i+third] for i in range(0,len(acts),third)][:3]
    while len(groups)<3: groups.append([])
    cards=[]
    for gi,g in enumerate(groups):
        cards.append((CARD_COLORS[gi], f"Labs {g[0]['num']}\u2013{g[-1]['num']}" if g else "\u2014",
                      [a["title"] for a in g] if g else ["\u2014"]))
    cards3(f"Hands-On Labs \u2014 {t['title']}", cards, kicker="WHAT YOU'LL DO")
    for ai,a in enumerate(acts):
        activity_overview(f"LAB {a['num']}", a["title"], a["desc"], a["build"], a["services"], kicker=f"TOPIC {t['code']} \u00b7 HANDS-ON")
        lab_steps(f"LAB {a['num']} \u00b7 STEPS", a["title"], a["steps"], per=6)
        test_slide(a["title"], a["test"], kicker=f"LAB {a['num']} \u00b7 VERIFY")
        if ai==len(acts)//2-1 and len(acts)>=4:
            brk("Lunch Break","1 hour")
    content(f"Recap \u2014 {t['title']}",
            ["You can now: "+a["objective"] for a in {x["objective"]:x for x in acts}.values()][:6],
            kicker="TOPIC RECAP", size=17)
    if tn in LAST_TOPIC_OF_DAY and LAST_TOPIC_OF_DAY[tn]<C.DAYS:
        big_statement(f"End of Day {LAST_TOPIC_OF_DAY[tn]}","Review today's labs and bring your questions to tomorrow morning.",f"DAY {LAST_TOPIC_OF_DAY[tn]} COMPLETE",color=TEAL)

# ---------------- CLOSE ----------------
section("WRAP-UP","Course Summary & Next Steps","")
tile_grid("What You Achieved",[(lo.split(":",1)[0], _lo_short(lo)) for lo in C.LEARNING_OUTCOMES],
 kicker="LEARNING OUTCOMES",cols=2,size=11)
tile_grid("Continuing Your Python Journey",[
 ("Rebuild the labs","Redo each lab from a blank folder until uv init, uv add and the pandas workflow are automatic."),
 ("Extend the capstone","Add your own data source, metric or page to the Streamlit dashboard you built."),
 ("Automate one real task","Pick a manual spreadsheet task in your own role and replace it with a documented script."),
 ("Keep the AI honest","Continue prompting AI tools for first drafts \u2014 and continue testing every line they produce.")],
 kicker="NEXT STEPS",cols=2,size=14)
content("Practice Exam",[
 "Sharpen your readiness with the Tertiary Infotech practice exams.",
 "Practice exams: https://exams.tertiaryinfotech.com",
 "Attempt the questions under timed conditions and review every explanation.",
 "Revisit any lab whose topic you miss, then re-take the practice questions."],kicker="TEST YOURSELF")
content("Assessment",[
 C.ASSESSMENT["written"], C.ASSESSMENT["practical"],
 f"Passing score: {C.ASSESSMENT['passing_score']}.",
 "Open book: slides, Learner Guide and lab files only.",
 "Remember to take the Assessment digital attendance (TRAQOM).",
 "Submit your completed answers on the LMS at https://lms-tms.tertiaryinfotech.com/."],kicker="WRAP-UP")
assessment_flow()
admin_attendance()
big_statement("Thank You!","You can now build, test and deploy AI-assisted Python programs for financial analysis \u2014 from a business requirement to a working Streamlit application.","WELL DONE",color=TEAL)

OUT=os.path.join(REPO,"courseware",f"{C.SHORT_TITLE}-{C.VERSION}.pptx")
prs.save(OUT)
print(f"Saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
